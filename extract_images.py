from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

prs = Presentation("10ADTW 2025 Potential Talents - 副本.pptx")

# Designer names extracted from slides
designers = ["ACT N°1", "AlienAnt", "ASSIGNMENTS", "AUBRUINO"]

# Create assets directory structure
base_path = Path("assets/designers")
base_path.mkdir(parents=True, exist_ok=True)

# Get slide dimensions (in EMUs - English Metric Units)
slide_width = prs.slide_width
slide_height = prs.slide_height

print(f"Slide dimensions: {slide_width} x {slide_height} EMUs")

# Extract images from each slide
image_count = 0
for slide_idx, slide in enumerate(prs.slides):
    designer_name = (
        designers[slide_idx]
        if slide_idx < len(designers)
        else f"Designer_{slide_idx + 1}"
    )

    # Create designer directory
    designer_path = base_path / designer_name
    designer_path.mkdir(parents=True, exist_ok=True)

    # Counter for images in this designer's slide
    image_num = 1

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Check if image is within slide bounds
            left = shape.left
            top = shape.top
            right = left + shape.width
            bottom = top + shape.height

            # Only extract if fully within slide bounds
            if (
                left >= 0
                and top >= 0
                and right <= slide_width
                and bottom <= slide_height
            ):
                image = shape.image
                image_bytes = image.blob

                # Determine file extension
                ext = image.ext
                filename = f"{image_num:02d}.{ext}"
                filepath = designer_path / filename

                # Save image
                with open(filepath, "wb") as f:
                    f.write(image_bytes)

                print(f"Saved: {filepath}")
                image_num += 1
                image_count += 1
            else:
                print(
                    f"Skipped image outside bounds in {designer_name}: ({left}, {top}) - ({right}, {bottom})"
                )

print(f"\nTotal images extracted: {image_count}")
